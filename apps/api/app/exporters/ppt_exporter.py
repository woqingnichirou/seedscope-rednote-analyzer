from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ..models import Note, Project
from .common import avg, ensure_parent, split_notes, top_counter


BLUE = RGBColor(31, 78, 121)
SLATE = RGBColor(68, 84, 106)
LIGHT = RGBColor(242, 246, 250)
WHITE = RGBColor(255, 255, 255)


def _add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.7))
    text_frame = box.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.size = Pt(28)
    paragraph.font.bold = True
    paragraph.font.color.rgb = BLUE
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.65), Inches(1.0), Inches(11.8), Inches(0.35))
        p = sub.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(12)
        p.font.color.rgb = SLATE


def _add_bullets(slide, bullets: list[str], top: float = 1.55) -> None:
    box = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.6), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18 if i == 0 else 16)
        p.font.color.rgb = SLATE
        p.space_after = Pt(10)


def _add_table(slide, headers: list[str], rows: list[list[str]], top: float = 1.5) -> None:
    row_count = len(rows) + 1
    col_count = len(headers)
    shape = slide.shapes.add_table(row_count, col_count, Inches(0.65), Inches(top), Inches(12.0), Inches(4.5))
    table = shape.table
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        p = cell.text_frame.paragraphs[0]
        p.text = header
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if r % 2 else WHITE
            p = cell.text_frame.paragraphs[0]
            p.text = str(value)
            p.font.size = Pt(11)
            p.font.color.rgb = SLATE


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def export_report_ppt(project: Project, notes: list[Note], path: Path) -> None:
    ensure_parent(path)
    brand_a_notes, brand_b_notes = split_notes(notes)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = _blank_slide(prs)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = LIGHT
    _add_title(slide, "SeedScope 竞品分析报告", f"{project.brand_a} vs {project.brand_b}｜{project.industry}｜{project.period}")
    _add_bullets(slide, ["本报告基于脱敏截图 OCR 与规则/模型标签生成。", "不包含真实品牌素材、真实截图或敏感预算。"], top=1.8)

    slide = _blank_slide(prs)
    _add_title(slide, "目录")
    _add_bullets(slide, ["核心结论", "样本概览", "标题与封面策略", "内容类型分布", "竞品优势与本品问题", "下阶段打法与行动清单"])

    slide = _blank_slide(prs)
    _add_title(slide, "核心结论", "每页只保留一个核心判断，便于内部汇报。")
    _add_bullets(slide, ["核心结论：Brand A 应优先提升决策型内容比例，而不是简单增加发布量。", "关键动作：补齐标题公式、封面证据位、正文判断标准和复投素材池。"])

    slide = _blank_slide(prs)
    _add_title(slide, "样本概览")
    _add_table(slide, ["指标", project.brand_a, project.brand_b], [["样本数", len(brand_a_notes), len(brand_b_notes)], ["平均点赞", f"{avg(brand_a_notes, 'likes'):.1f}", f"{avg(brand_b_notes, 'likes'):.1f}"], ["平均收藏", f"{avg(brand_a_notes, 'collects'):.1f}", f"{avg(brand_b_notes, 'collects'):.1f}"], ["平均评论", f"{avg(brand_a_notes, 'comments'):.1f}", f"{avg(brand_b_notes, 'comments'):.1f}"]])

    slide = _blank_slide(prs)
    _add_title(slide, "标题策略对比")
    _add_table(slide, ["标题类型", "样本数"], [[k, v] for k, v in top_counter(notes, "title_type", 6)])

    slide = _blank_slide(prs)
    _add_title(slide, "封面策略对比")
    _add_table(slide, ["封面类型", "样本数"], [[k, v] for k, v in top_counter(notes, "cover_type", 6)])

    slide = _blank_slide(prs)
    _add_title(slide, "内容类型分布")
    _add_table(slide, ["内容类型", "样本数"], [[k, v] for k, v in top_counter(notes, "content_type", 6)])

    slide = _blank_slide(prs)
    _add_title(slide, "竞品优势")
    _add_bullets(slide, ["竞品优势通常来自更强的问题入口、方法论表达和结果证据。", "建议对标其内容结构，而不是复制单篇表达。"])

    slide = _blank_slide(prs)
    _add_title(slide, "本品问题")
    _add_bullets(slide, ["本品问题：品牌前置偏多、证据位不足、方法论段落偏弱、CTA 偏硬。", "优先级：先改内容结构，再扩大投放。"])

    slide = _blank_slide(prs)
    _add_title(slide, "下阶段打法")
    _add_bullets(slide, ["第 1-2 周：模板重建。", "第 3-4 周：小规模赛马。", "第 5-8 周：白名单复投。", "第 9-12 周：品效联动。"])

    slide = _blank_slide(prs)
    _add_title(slide, "行动清单")
    _add_bullets(slide, ["建立达人能力标签。", "输出标题公式库和封面模板库。", "固化正文 SOP。", "定义素材入池阈值。", "建立周度复盘看板。"])

    for slide in prs.slides:
        footer = slide.shapes.add_textbox(Inches(0.65), Inches(7.05), Inches(12), Inches(0.25))
        p = footer.text_frame.paragraphs[0]
        p.text = "SeedScope | Local-first, screenshot-based, brand-safe analysis"
        p.font.size = Pt(9)
        p.font.color.rgb = SLATE
        p.alignment = PP_ALIGN.RIGHT

    prs.save(path)
